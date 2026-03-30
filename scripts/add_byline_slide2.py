#!/usr/bin/env python3
"""
add_byline_slide2.py
Adds a "Built by Dana" byline pill to the top-left corner of slide 2
in newsletter_journey.pptx, matching the HTML orbit-frame style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH  = (r"c:\Users\danaa\OneDrive - AMDOCS\Backup Folders\Documents"
              r"\CC Projects\AI newsletter\newsletter_journey.pptx")
PHOTO_PATH = (r"c:\Users\danaa\OneDrive - AMDOCS\Backup Folders\Documents"
              r"\CC Projects\AI newsletter\Image (13).jpg")

BLUE  = "00D4FF"
TEXT  = "E8E8F0"
MUTED = "6B6B80"
DARK  = "0D0D1C"

def rgb(h):
    return RGBColor.from_string(h)

def make_pill_shape(slide, x, y, w, h):
    """Rounded rectangle with max corner radius = pill / stadium shape."""
    s = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(DARK)
    s.line.color.rgb = rgb("252535")
    s.line.width = Pt(1.0)
    # Max corner rounding → pill shape
    sp_el = s._element
    prstGeom = sp_el.find('.//' + qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        for child in list(avLst):
            avLst.remove(child)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 50000')   # 50% = full pill
    return s

def add_circle(slide, x, y, d, fill=None, border=None, border_pt=1.5):
    """Add a circle (oval) shape."""
    s = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(d), Inches(d))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = rgb(border)
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s

def add_photo_circle(slide, photo_path, x, y, d):
    """Add a picture clipped to a circle."""
    pic = slide.shapes.add_picture(photo_path, Inches(x), Inches(y), Inches(d), Inches(d))
    # Clip to ellipse by changing prstGeom in spPr
    sp_el = pic._element
    spPr  = sp_el.find(qn('p:spPr'))
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'ellipse')
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is not None:
            for child in list(avLst):
                avLst.remove(child)
    # Crop photo to top portion (face focus)
    blipFill = sp_el.find(qn('p:blipFill'))
    if blipFill is not None:
        srcRect = etree.Element(qn('a:srcRect'))
        srcRect.set('b', '25000')   # trim 25% from bottom
        blipFill.insert(0, srcRect)
    return pic

def add_txt(slide, text, x, y, w, h, size=12, bold=False,
            color=TEXT, align=PP_ALIGN.LEFT, font="Arial"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = rgb(color)
    run.font.name      = font
    return txb

# ── OPEN PPTX & TARGET SLIDE 2 ───────────────────────────────────────────────
prs   = Presentation(PPTX_PATH)
slide = prs.slides[1]   # index 1 = slide 2

# ── LAYOUT ───────────────────────────────────────────────────────────────────
PILL_X = 0.16
PILL_Y = 0.13
PILL_W = 2.08
PILL_H = 0.64

PHOTO_D  = 0.46                                          # photo diameter
PHOTO_X  = PILL_X + 0.09
PHOTO_Y  = PILL_Y + (PILL_H - PHOTO_D) / 2

RING_D   = PHOTO_D + 0.07                                # ring slightly larger
RING_X   = PHOTO_X - 0.035
RING_Y   = PHOTO_Y - 0.035

BALL_D   = 0.075                                         # orbit dot
BALL_X   = RING_X + RING_D / 2 - BALL_D / 2
BALL_Y   = RING_Y - BALL_D / 2

TEXT_X   = PHOTO_X + PHOTO_D + 0.12
LABEL_Y  = PILL_Y + 0.08
NAME_Y   = PILL_Y + 0.30

# ── DRAW ELEMENTS (back to front) ────────────────────────────────────────────

# 1. Pill background
make_pill_shape(slide, PILL_X, PILL_Y, PILL_W, PILL_H)

# 2. Photo clipped to circle
add_photo_circle(slide, PHOTO_PATH, PHOTO_X, PHOTO_Y, PHOTO_D)

# 3. Neon blue ring around photo
ring = add_circle(slide, RING_X, RING_Y, RING_D, fill=None, border=BLUE, border_pt=1.5)

# 4. Orbit dot (glowing ball at top of ring)
add_circle(slide, BALL_X, BALL_Y, BALL_D, fill=BLUE, border=None)

# 5. "CURATED BY" label
add_txt(slide, "CURATED BY",
        x=TEXT_X, y=LABEL_Y, w=1.2, h=0.20,
        size=7.5, bold=True, color=MUTED, align=PP_ALIGN.LEFT)

# 6. Name
add_txt(slide, "Dana Aronovich",
        x=TEXT_X, y=NAME_Y, w=1.3, h=0.26,
        size=13, bold=True, color=TEXT, align=PP_ALIGN.LEFT)

# ── SAVE ─────────────────────────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f"Done — byline added to slide 2 in:\n{PPTX_PATH}")

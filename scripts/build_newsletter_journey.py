#!/usr/bin/env python3
"""
build_newsletter_journey.pptx
Single slide: "Building the AI Newsletter" journey
Dark-tech style matching the AI Newsflash newsletter.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── PALETTE (matches newsletter CSS variables) ────────────────────────────────
BG      = "0A0A0F"
SURFACE = "0F0F1C"
DIM     = "1A1A2E"
TEXT    = "E8E8F0"
MUTED   = "6B6B80"
BLUE    = "00D4FF"
PINK    = "FF006E"
GREEN   = "00FF9D"
YELLOW  = "FFD60A"

STEP_COLORS = [YELLOW, BLUE, GREEN, PINK, "00D4FF", "00FF9D"]

STEPS = [
    ("💡", "0", "Define the Vision",  ["Prompted ChatGPT to", "shape the concept"]),
    ("🔍", "1", "Curate Content",     ["Select 10–11 top", "AI stories weekly"]),
    ("✍",  "2", "Build the HTML",    ["Dark-tech template", "& card structure"]),
    ("🎨", "3", "AI Hero Images",     ["Imagen 4 generates", "cinematic visuals"]),
    ("🌐", "4", "Publish Online",     ["GitHub Pages", "live public URL"]),
    ("✉",  "5", "Send by Email",     ["Gmail SMTP to", "all subscribers"]),
]

# ── HELPERS ───────────────────────────────────────────────────────────────────

def rgb(h):
    return RGBColor.from_string(h)


def add_rect(slide, x, y, w, h, fill=SURFACE, border=None, border_pt=1.5):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    if border:
        s.line.color.rgb = rgb(border)
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s


def add_txt(slide, text, x, y, w, h, size=14, bold=False,
            color=TEXT, align=PP_ALIGN.LEFT, font="Arial"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.name = font
    return txb


def add_lines(slide, lines, x, y, w, h, size=11, color=MUTED,
              align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
        run.font.name = "Arial"
    return txb


# ── PRESENTATION ──────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

# Background
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = rgb(BG)

# ── TOP ACCENT BAR (gradient simulation) ─────────────────────────────────────
bar_colors = [PINK, BLUE, GREEN, YELLOW, BLUE, PINK]
bw = 13.33 / len(bar_colors)
for i, c in enumerate(bar_colors):
    add_rect(slide, i * bw, 0, bw + 0.02, 0.07, fill=c)

# ── TITLE & SUBTITLE ──────────────────────────────────────────────────────────
add_txt(slide, "🚀  Building the AI Newsletter",
        x=0.5, y=0.12, w=10.5, h=0.65,
        size=30, bold=True, color=TEXT, align=PP_ALIGN.LEFT,
        font="Segoe UI Emoji")

add_txt(slide, "From idea to inbox — a fully AI-powered workflow",
        x=0.5, y=0.78, w=10.5, h=0.38,
        size=13, color=MUTED, align=PP_ALIGN.LEFT)

# ── LAYOUT CONSTANTS ──────────────────────────────────────────────────────────
BW   = 3.55   # box width
BH   = 2.52   # box height
GAP  = 0.43   # gap between boxes (holds arrow)
LM   = 0.90   # left margin
R1Y  = 1.46   # row 1 top y
R2Y  = 4.28   # row 2 top y

# 3 column x-positions
COL = [LM, LM + BW + GAP, LM + 2 * (BW + GAP)]

# Snake layout: row1 L→R (steps 0,1,2), row2 R→L (steps 3,4,5)
POS = {
    0: (COL[0], R1Y),
    1: (COL[1], R1Y),
    2: (COL[2], R1Y),
    3: (COL[2], R2Y),   # same column as step 2
    4: (COL[1], R2Y),
    5: (COL[0], R2Y),
}

# ── STEP BOXES ────────────────────────────────────────────────────────────────

for si in range(6):
    emoji, num, title, desc = STEPS[si]
    accent = STEP_COLORS[si]
    bx, by = POS[si]

    # ── Box body
    add_rect(slide, bx, by, BW, BH, fill=SURFACE, border=accent, border_pt=1.2)

    # ── Top color strip
    add_rect(slide, bx, by, BW, 0.08, fill=accent)

    # ── Step number badge (small filled square, top-left)
    add_rect(slide, bx + 0.13, by + 0.16, 0.30, 0.28, fill=accent)
    add_txt(slide, num,
            x=bx + 0.13, y=by + 0.15, w=0.30, h=0.28,
            size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)

    # ── Emoji (large, centered)
    add_txt(slide, emoji,
            x=bx + 0.1, y=by + 0.12, w=BW - 0.2, h=0.68,
            size=30, color=TEXT, align=PP_ALIGN.CENTER,
            font="Segoe UI Emoji")

    # ── Title
    add_txt(slide, title,
            x=bx + 0.1, y=by + 0.88, w=BW - 0.2, h=0.52,
            size=15, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    # ── Accent underline under title
    add_rect(slide, bx + BW/2 - 0.4, by + 1.44, 0.80, 0.03, fill=accent)

    # ── Description lines
    add_lines(slide, desc,
              x=bx + 0.15, y=by + 1.54, w=BW - 0.3, h=0.82,
              size=11.5, color=MUTED, align=PP_ALIGN.CENTER)

# ── ARROWS ────────────────────────────────────────────────────────────────────

ARROW_COLOR = "3A3A5E"
AH = 0.28

# Row 1 right arrows (steps 0→1 and 1→2)
AY1 = R1Y + BH / 2 - AH / 2
for i in range(2):
    ax = COL[i] + BW + 0.04
    aw = GAP - 0.08
    add_txt(slide, "▶",
            x=ax, y=AY1, w=aw, h=AH,
            size=20, color=ARROW_COLOR, align=PP_ALIGN.CENTER)

# Down arrow connector (step 2 → step 3, same column)
cx = COL[2] + BW / 2 - 0.2
cy = R1Y + BH + 0.06
ch = R2Y - R1Y - BH - 0.06
add_txt(slide, "▼",
        x=cx, y=cy, w=0.4, h=ch,
        size=20, color=ARROW_COLOR, align=PP_ALIGN.CENTER)

# Row 2 left arrows (steps 3←4 and 4←5)
# Step 3 at COL[2], step 4 at COL[1]: arrow in COL[1]+BW gap
# Step 4 at COL[1], step 5 at COL[0]: arrow in COL[0]+BW gap
AY2 = R2Y + BH / 2 - AH / 2
for i in range(2):
    ax = COL[i] + BW + 0.04
    aw = GAP - 0.08
    add_txt(slide, "◀",
            x=ax, y=AY2, w=aw, h=AH,
            size=20, color=ARROW_COLOR, align=PP_ALIGN.CENTER)

# ── FLOW LABEL ────────────────────────────────────────────────────────────────
# Small label above row 1 arrows showing direction
add_txt(slide, "workflow",
        x=0.5, y=R1Y - 0.28, w=1.2, h=0.25,
        size=8, color=MUTED, align=PP_ALIGN.LEFT)

# ── FOOTER ────────────────────────────────────────────────────────────────────
add_rect(slide, 0.5, 7.06, 12.33, 0.015, fill=DIM)
add_txt(slide,
        "Built with  ChatGPT  ·  Claude Code  ·  Imagen 4 (Google)  ·  GitHub Pages  ·  Gmail SMTP",
        x=0.5, y=7.1, w=12.33, h=0.32,
        size=9.5, color=MUTED, align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────────────────────────
out = (
    r"c:\Users\danaa\OneDrive - AMDOCS\Backup Folders\Documents"
    r"\CC Projects\AI newsletter\newsletter_journey.pptx"
)
prs.save(out)
print(f"Saved: {out}")
